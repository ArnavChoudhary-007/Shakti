"""
tests/create_samples.py
Creates small sample files for each connector type so Phase 1 tests
can run without needing real user files.
Run: python -m tests.create_samples
"""
from __future__ import annotations

import json
import os
from pathlib import Path

SAMPLES_DIR = Path(__file__).parent / "sample_files"
SAMPLES_DIR.mkdir(exist_ok=True)


def create_pdf_sample() -> None:
    """Create a minimal digital PDF using pymupdf."""
    try:
        import fitz
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((72, 100), "RAG Pipeline Test Document\n\nThis is page 1 of a sample PDF.\nIt contains some text for testing the PDF connector.")
        page2 = doc.new_page()
        page2.insert_text((72, 100), "Page 2\n\nMore content on the second page.\nThis tests multi-page citation metadata.")
        doc.save(str(SAMPLES_DIR / "sample.pdf"))
        doc.close()
        print("  [OK] Created sample.pdf")
    except Exception as e:
        print(f"  [SKIP] PDF sample: {e}")


def create_docx_sample() -> None:
    try:
        from docx import Document
        doc = Document()
        doc.add_heading("Sample Word Document", level=1)
        doc.add_paragraph("This is the first paragraph of a sample Word document.")
        doc.add_paragraph("This is the second paragraph, used to test DOCX connector extraction.")
        t = doc.add_table(rows=2, cols=3)
        t.rows[0].cells[0].text = "Name"
        t.rows[0].cells[1].text = "Amount"
        t.rows[0].cells[2].text = "Date"
        t.rows[1].cells[0].text = "Acme Corp"
        t.rows[1].cells[1].text = "5000"
        t.rows[1].cells[2].text = "2024-01-15"
        doc.save(str(SAMPLES_DIR / "sample.docx"))
        print("  [OK] Created sample.docx")
    except Exception as e:
        print(f"  [SKIP] DOCX sample: {e}")


def create_pptx_sample() -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Sample Presentation Slide 1"
        slide.placeholders[1].text = "This is the content of slide 1.\nBullet point one.\nBullet point two."
        slide2 = prs.slides.add_slide(slide_layout)
        slide2.shapes.title.text = "Sample Presentation Slide 2"
        slide2.placeholders[1].text = "Second slide content for testing PPTX connector."
        prs.save(str(SAMPLES_DIR / "sample.pptx"))
        print("  [OK] Created sample.pptx")
    except Exception as e:
        print(f"  [SKIP] PPTX sample: {e}")


def create_excel_sample() -> None:
    try:
        import pandas as pd
        # Sheet 1: tabular (ledger)
        df_ledger = pd.DataFrame({
            "Vendor": ["Acme Corp", "Beta Ltd", "Gamma Inc"],
            "Amount": [5000.00, 12500.50, 750.00],
            "Date": ["2024-01-15", "2024-01-20", "2024-01-25"],
            "Invoice_No": ["INV-001", "INV-002", "INV-003"],
            "Status": ["Paid", "Pending", "Paid"],
        })
        # Sheet 2: narrative (text-heavy)
        df_notes = pd.DataFrame({
            "Note": [
                "Q1 vendor meeting minutes: discussed payment terms",
                "Follow-up required on Beta Ltd outstanding balance",
                "Gamma Inc payment confirmed via bank transfer",
            ],
        })
        with pd.ExcelWriter(str(SAMPLES_DIR / "sample.xlsx"), engine="openpyxl") as writer:
            df_ledger.to_excel(writer, sheet_name="Payments", index=False)
            df_notes.to_excel(writer, sheet_name="Notes", index=False)
        print("  [OK] Created sample.xlsx (2 sheets: Payments=tabular, Notes=narrative)")
    except Exception as e:
        print(f"  [SKIP] Excel sample: {e}")


def create_csv_sample() -> None:
    try:
        content = "Vendor,Amount,Date,Invoice_No\nAcme Corp,5000.00,2024-01-15,INV-001\nBeta Ltd,12500.50,2024-01-20,INV-002\n"
        (SAMPLES_DIR / "sample.csv").write_text(content, encoding="utf-8")
        print("  [OK] Created sample.csv")
    except Exception as e:
        print(f"  [SKIP] CSV sample: {e}")


def create_mbox_sample() -> None:
    try:
        import mailbox, email.message
        mbox = mailbox.mbox(str(SAMPLES_DIR / "sample.mbox"))
        msg = email.message.Message()
        msg["From"] = "alice@example.com"
        msg["To"] = "bob@example.com"
        msg["Subject"] = "Test email for RAG connector"
        msg["Date"] = "Mon, 15 Jan 2024 10:00:00 +0000"
        msg["Message-ID"] = "<test-001@example.com>"
        msg.set_payload("Hello Bob,\n\nThis is a test email for the RAG pipeline email connector.\nIt contains some content to verify citation metadata.\n\nBest,\nAlice")
        mbox.add(msg)
        mbox.close()
        print("  [OK] Created sample.mbox")
    except Exception as e:
        print(f"  [SKIP] mbox sample: {e}")


def create_whatsapp_sample() -> None:
    content = (
        "1/15/2024, 10:00 AM - Alice: Hello Bob! Did you see the invoice from Acme?\n"
        "1/15/2024, 10:01 AM - Bob: Not yet, let me check.\n"
        "1/15/2024, 10:02 AM - Bob: Found it. The total is $5,000 due by end of month.\n"
        "1/15/2024, 10:03 AM - Alice: Perfect, I'll process the payment today.\n"
        "1/15/2024, 10:05 AM - Alice: Done! Payment sent via bank transfer.\n"
    )
    (SAMPLES_DIR / "sample_whatsapp.txt").write_text(content, encoding="utf-8")
    print("  [OK] Created sample_whatsapp.txt")


def create_telegram_sample() -> None:
    data = {
        "name": "Project Finance Chat",
        "type": "private_group",
        "messages": [
            {
                "id": 1,
                "type": "message",
                "date": "2024-01-15T10:00:00",
                "from": "Alice",
                "from_id": "user123",
                "text": "Has anyone reviewed the Q4 invoices?",
            },
            {
                "id": 2,
                "type": "message",
                "date": "2024-01-15T10:01:00",
                "from": "Bob",
                "from_id": "user456",
                "text": [
                    "Yes, I checked. Acme Corp invoice ",
                    {"type": "bold", "text": "INV-001"},
                    " is $5,000 due Jan 31.",
                ],
            },
            {
                "id": 3,
                "type": "service",
                "date": "2024-01-15T10:02:00",
                "actor": "Bob",
                "action": "pin_message",
                "text": "",
            },
        ],
    }
    (SAMPLES_DIR / "sample_telegram.json").write_text(
        json.dumps(data, indent=2), encoding="utf-8"
    )
    print("  [OK] Created sample_telegram.json")


def create_slack_sample() -> None:
    slack_dir = SAMPLES_DIR / "slack_export" / "general"
    slack_dir.mkdir(parents=True, exist_ok=True)

    users = [
        {"id": "U001", "name": "alice", "real_name": "Alice Smith", "profile": {"display_name": "alice"}},
        {"id": "U002", "name": "bob", "real_name": "Bob Jones", "profile": {"display_name": "bob"}},
    ]
    (SAMPLES_DIR / "slack_export" / "users.json").write_text(
        json.dumps(users, indent=2), encoding="utf-8"
    )

    channels = [{"id": "C001", "name": "general", "topic": {"value": "General discussions"}, "purpose": {"value": "Team channel"}}]
    (SAMPLES_DIR / "slack_export" / "channels.json").write_text(
        json.dumps(channels, indent=2), encoding="utf-8"
    )

    messages = [
        {"type": "message", "user": "U001", "text": "Hey team, the invoice from Acme <@U002> has been approved.", "ts": "1705312800.000100"},
        {"type": "message", "user": "U002", "text": "Great! I will process the payment by EOD.", "ts": "1705312860.000200"},
        {"type": "message", "subtype": "channel_join", "user": "U003", "text": "joined", "ts": "1705312900.000300"},
    ]
    (slack_dir / "2024-01-15.json").write_text(
        json.dumps(messages, indent=2), encoding="utf-8"
    )
    print("  [OK] Created slack_export/ directory")


def create_teams_sample() -> None:
    data = {
        "TeamName": "Finance Team",
        "ChannelName": "General",
        "Messages": [
            {
                "Id": "msg-001",
                "CreatedDateTime": "2024-01-15T10:00:00Z",
                "From": {"User": {"DisplayName": "Alice Smith"}},
                "Body": {"Content": "The Q4 vendor payments have been reconciled.", "ContentType": "text"},
                "Replies": [
                    {
                        "Id": "msg-001-r1",
                        "CreatedDateTime": "2024-01-15T10:05:00Z",
                        "From": {"User": {"DisplayName": "Bob Jones"}},
                        "Body": {"Content": "<p>Thanks Alice, I will update the ledger.</p>", "ContentType": "html"},
                        "Replies": [],
                    }
                ],
            }
        ],
    }
    (SAMPLES_DIR / "sample_teams.json").write_text(
        json.dumps(data, indent=2), encoding="utf-8"
    )
    print("  [OK] Created sample_teams.json")


def create_invoice_sample() -> None:
    try:
        import fitz
        doc = fitz.open()
        page = doc.new_page()
        invoice_text = """INVOICE

Invoice No: INV-2024-0042
Invoice Date: January 15, 2024
Due Date: February 15, 2024

From:
Acme Corporation
123 Business Street, Suite 100
New York, NY 10001

Bill To:
Beta Technologies Ltd
456 Client Avenue
San Francisco, CA 94102

Description                 Qty    Unit Price    Amount
---------------------------------------------------------
Software License (Annual)    1     $4,000.00     $4,000.00
Support & Maintenance        1       $750.00       $750.00
Training Sessions            2       $125.00       $250.00
---------------------------------------------------------
                                    TOTAL:       $5,000.00

Payment Terms: Net 30
Bank Transfer: Account 1234567890, Routing 021000021
"""
        page.insert_text((50, 50), invoice_text, fontsize=10)
        doc.save(str(SAMPLES_DIR / "sample_invoice.pdf"))
        doc.close()
        print("  [OK] Created sample_invoice.pdf")
    except Exception as e:
        print(f"  [SKIP] Invoice PDF sample: {e}")


if __name__ == "__main__":
    print(f"Creating sample files in: {SAMPLES_DIR}")
    create_pdf_sample()
    create_docx_sample()
    create_pptx_sample()
    create_excel_sample()
    create_csv_sample()
    create_mbox_sample()
    create_whatsapp_sample()
    create_telegram_sample()
    create_slack_sample()
    create_teams_sample()
    create_invoice_sample()
    print("\nDone. Run: python -m tests.test_connectors")
