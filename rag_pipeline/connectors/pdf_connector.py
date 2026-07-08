"""
connectors/pdf_connector.py
Handles: PDF (.pdf), Word (.docx/.doc), PowerPoint (.pptx/.ppt)

Strategy:
  PDF  — pymupdf for digital text; page-level OCR (pytesseract) if a page
          has no text layer (scanned). Emits one envelope per page for
          fine-grained citation_meta.
  DOCX — python-docx; extracts paragraphs + table cells.
  PPTX — python-pptx; extracts all text frames per slide.
  DOC / PPT (legacy) — not handled natively; user advised to convert first.

Each envelope represents ONE PAGE (PDF/PPTX) or the whole document (DOCX),
so citation_meta can be page-specific.
"""
from __future__ import annotations

import io
import logging
import os
from pathlib import Path
from typing import Any, Dict, Generator, List, Optional

from .base import envelope, make_source_id

logger = logging.getLogger(__name__)

# ── Optional imports (graceful degradation) ─────────────────
try:
    import fitz  # pymupdf
    _HAS_FITZ = True
except ImportError:
    _HAS_FITZ = False
    logger.warning("pymupdf not installed — PDF support disabled.")

try:
    from PIL import Image
    import pytesseract
    _HAS_OCR = True
except ImportError:
    _HAS_OCR = False
    logger.warning("Pillow/pytesseract not installed — OCR fallback disabled.")

try:
    from docx import Document as DocxDocument
    _HAS_DOCX = True
except ImportError:
    _HAS_DOCX = False
    logger.warning("python-docx not installed — DOCX support disabled.")

try:
    from pptx import Presentation
    from pptx.util import Pt
    _HAS_PPTX = True
except ImportError:
    _HAS_PPTX = False
    logger.warning("python-pptx not installed — PPTX support disabled.")


# ─────────────────────────────────────────────────────────────
# PDF
# ─────────────────────────────────────────────────────────────

def _ocr_page(page: Any) -> str:
    """Render a fitz page to an image and run Tesseract OCR."""
    if not _HAS_OCR:
        return ""
    mat = fitz.Matrix(2.0, 2.0)          # 2x zoom for better OCR accuracy
    pix = page.get_pixmap(matrix=mat)
    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
    return pytesseract.image_to_string(img)


def _is_scanned_page(page: Any, min_chars: int = 30) -> bool:
    """Return True if a page has fewer than min_chars of extractable text."""
    return len(page.get_text("text").strip()) < min_chars


def extract_pdf(
    file_path: str,
    use_ocr: bool = True,
) -> Generator[Dict[str, Any], None, None]:
    """
    Yield one envelope per PDF page.
    Automatically falls back to OCR for pages with no text layer.
    """
    if not _HAS_FITZ:
        raise RuntimeError("pymupdf is required for PDF extraction.")

    path = Path(file_path)
    source_id = make_source_id(file_path)
    file_name = path.name

    doc = fitz.open(file_path)
    total_pages = len(doc)

    for page_num, page in enumerate(doc, start=1):
        text = page.get_text("text").strip()
        used_ocr = False

        if use_ocr and _is_scanned_page(page):
            text = _ocr_page(page).strip()
            used_ocr = True

        if not text:
            logger.debug("Page %d/%d has no text — skipping.", page_num, total_pages)
            continue

        # Extract images on page for metadata
        image_count = len(page.get_images(full=True))

        yield envelope(
            source_type="pdf",
            source_id=f"{source_id}_p{page_num}",
            raw_path=file_path,
            raw_metadata={
                "file_name": file_name,
                "page_number": page_num,
                "total_pages": total_pages,
                "used_ocr": used_ocr,
                "image_count": image_count,
                "title": doc.metadata.get("title", "") or file_name,
                "author": doc.metadata.get("author", ""),
                "created": doc.metadata.get("creationDate", ""),
                "full_text": text,
            },
            modality="text",
        )

    doc.close()


# ─────────────────────────────────────────────────────────────
# DOCX
# ─────────────────────────────────────────────────────────────

def _extract_docx_text(doc_obj: Any) -> str:
    """Extract paragraphs + table cells from a DOCX document object."""
    parts: List[str] = []

    # Paragraphs (preserves order via XML body iteration)
    for para in doc_obj.paragraphs:
        stripped = para.text.strip()
        if stripped:
            parts.append(stripped)

    # Table cells
    for table in doc_obj.tables:
        for row in table.rows:
            row_text = " | ".join(
                cell.text.strip() for cell in row.cells if cell.text.strip()
            )
            if row_text:
                parts.append(row_text)

    return "\n\n".join(parts)


def extract_docx(file_path: str) -> Generator[Dict[str, Any], None, None]:
    """
    Yield a single envelope for a DOCX file.
    (Word docs don't have stable page boundaries in python-docx,
    so we treat the whole doc as one unit; chunker will split it.)
    """
    if not _HAS_DOCX:
        raise RuntimeError("python-docx is required for DOCX extraction.")

    path = Path(file_path)
    source_id = make_source_id(file_path)
    doc_obj = DocxDocument(file_path)

    # Try to get core properties
    props = doc_obj.core_properties
    title = props.title or path.stem
    author = props.author or ""
    created = str(props.created) if props.created else ""

    text = _extract_docx_text(doc_obj)

    if not text.strip():
        logger.warning("DOCX %s appears to have no extractable text.", path.name)
        return

    yield envelope(
        source_type="docx",
        source_id=source_id,
        raw_path=file_path,
        raw_metadata={
            "file_name": path.name,
            "title": title,
            "author": author,
            "created": created,
            "paragraph_count": len(doc_obj.paragraphs),
            "table_count": len(doc_obj.tables),
            "full_text": text,          # normalizer reads this
        },
        modality="text",
    )


# ─────────────────────────────────────────────────────────────
# PPTX
# ─────────────────────────────────────────────────────────────

def _extract_slide_text(slide: Any) -> str:
    """Extract all text frames from a slide, preserving order."""
    parts: List[str] = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            line = " ".join(run.text for run in para.runs).strip()
            if line:
                parts.append(line)
    return "\n".join(parts)


def extract_pptx(file_path: str) -> Generator[Dict[str, Any], None, None]:
    """Yield one envelope per PPTX slide."""
    if not _HAS_PPTX:
        raise RuntimeError("python-pptx is required for PPTX extraction.")

    path = Path(file_path)
    source_id = make_source_id(file_path)
    prs = Presentation(file_path)
    total_slides = len(prs.slides)

    # Core properties
    props = prs.core_properties
    title = props.title or path.stem
    author = props.author or ""
    created = str(props.created) if props.created else ""

    for slide_num, slide in enumerate(prs.slides, start=1):
        text = _extract_slide_text(slide)
        if not text.strip():
            logger.debug("Slide %d/%d has no text — skipping.", slide_num, total_slides)
            continue

        # Slide title (first placeholder of type TITLE or CENTER_TITLE)
        slide_title = ""
        try:
            slide_title = slide.shapes.title.text.strip() if slide.shapes.title else ""
        except Exception:
            pass

        yield envelope(
            source_type="pptx",
            source_id=f"{source_id}_s{slide_num}",
            raw_path=file_path,
            raw_metadata={
                "file_name": path.name,
                "slide_number": slide_num,
                "total_slides": total_slides,
                "slide_title": slide_title,
                "presentation_title": title,
                "author": author,
                "created": created,
                "full_text": text,
            },
            modality="text",
        )


# ─────────────────────────────────────────────────────────────
# Unified entry point
# ─────────────────────────────────────────────────────────────

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".doc", ".pptx", ".ppt"}


def extract(
    file_path: str,
    use_ocr: bool = True,
) -> Generator[Dict[str, Any], None, None]:
    """
    Route to the correct extractor based on file extension.
    Yields one or more connector envelopes.
    """
    ext = Path(file_path).suffix.lower()

    if ext == ".pdf":
        yield from extract_pdf(file_path, use_ocr=use_ocr)
    elif ext in (".docx", ".doc"):
        yield from extract_docx(file_path)
    elif ext in (".pptx", ".ppt"):
        yield from extract_pptx(file_path)
    else:
        raise ValueError(
            f"Unsupported extension {ext!r}. "
            f"Supported: {sorted(SUPPORTED_EXTENSIONS)}"
        )
